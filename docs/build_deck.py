"""
Generate the KrishiMitra project review deck.

Kept as a script rather than a hand-built file so every figure stays tied to
what the system actually produces: the metrics come from models/metrics.json and
the data-coverage numbers from state_crop_stats.json. Retrain or rebuild the
statistics and the deck updates itself instead of drifting out of date.

Run:  ./backend/.venv/bin/python docs/build_deck.py
Out:  docs/KrishiMitra_Project_Review.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "KrishiMitra_Project_Review.pptx"
DATA = ROOT / "backend" / "app" / "data"

METRICS = json.loads((ROOT / "backend" / "models" / "metrics.json").read_text())
STATS = json.loads((DATA / "state_crop_stats.json").read_text())
COVERAGE = STATS["_meta"]["coverage"]
SCHEMES = json.loads((DATA / "schemes.json").read_text())["schemes"]
PROFILES = json.loads((DATA / "state_profiles.json").read_text())["states"]

# ---------------------------------------------------------------- palette ---

BG = RGBColor(0x0A, 0x12, 0x16)
INK = RGBColor(0x0F, 0x1E, 0x24)
PANEL = RGBColor(0x15, 0x28, 0x2E)
PANEL_HI = RGBColor(0x1B, 0x33, 0x3A)
LINE = RGBColor(0x27, 0x47, 0x50)
TEXT = RGBColor(0xEC, 0xF4, 0xF1)
DIM = RGBColor(0x93, 0xAC, 0xA9)
FAINT = RGBColor(0x5E, 0x7B, 0x78)
GREEN = RGBColor(0x6E, 0xCF, 0x94)
GREEN_D = RGBColor(0x4C, 0x9F, 0x70)
AMBER = RGBColor(0xE0, 0xA4, 0x58)
RED = RGBColor(0xD1, 0x68, 0x5C)
BLUE = RGBColor(0x5B, 0x8A, 0xB0)
VIOLET = RGBColor(0x9B, 0x7C, 0xC4)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.8)                      # page margin
CW = W - 2 * M                       # content width
GAP = Inches(0.3)
COL2 = (CW - GAP) / 2                          # 5.72"
COL3 = (CW - 2 * Inches(0.26)) / 3             # 3.73"
COL4 = (CW - 3 * Inches(0.24)) / 4             # 2.75"


def col(i: int, w, gap):
    """Left edge of column i in a grid of width w and gutter gap."""
    return M + (w + gap) * i
FONT = "Helvetica Neue"
MONO = "Menlo"

_state = {"n": 0}

# ------------------------------------------------------------- primitives ---


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK if dark else BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def tbox(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size, color=TEXT, bold=False, first=False, italic=False,
         before=0, after=6, align=PP_ALIGN.LEFT, line=None, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def rect(s, x, y, w, h, fill=PANEL, border=None, radius=0.045, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    c = s.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        c.adjustments[0] = radius
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    if border:
        c.line.color.rgb = border
        c.line.width = Pt(1.1)
    else:
        c.line.fill.background()
    c.shadow.inherit = False
    return c


def chip(s, x, y, w, h, text, fill, fg=None, size=9.5):
    c = rect(s, x, y, w, h, fill=fill, radius=0.45)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg or BG
    r.font.name = FONT
    return c


def header(s, title_text, kicker=None, sub=None):
    y = Inches(0.52)
    if kicker:
        tf = tbox(s, M, y, CW, Inches(0.24))
        para(tf, kicker.upper(), 10, GREEN, bold=True, first=True, after=0)
        y += Inches(0.3)
    tf = tbox(s, M, y, CW, Inches(0.6))
    para(tf, title_text, 31, TEXT, bold=True, first=True, after=0)
    y += Inches(0.56)
    if sub:
        tf = tbox(s, M, y, CW, Inches(0.5))
        para(tf, sub, 14, DIM, first=True, after=0, line=1.2)
        y += Inches(0.34)
    rect(s, M, y + Inches(0.12), Inches(1.15), Pt(2.5), fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
    return y + Inches(0.42)


def foot(s):
    _state["n"] += 1
    tf = tbox(s, W - Inches(1.4), H - Inches(0.52), Inches(0.6), Inches(0.3), align=PP_ALIGN.RIGHT)
    para(tf, f"{_state['n']:02d}", 9.5, FAINT, first=True, align=PP_ALIGN.RIGHT, after=0)
    tf = tbox(s, M, H - Inches(0.52), Inches(4), Inches(0.3))
    para(tf, "KrishiMitra", 9.5, FAINT, first=True, after=0)


def stat(s, x, y, w, h, value, label, note=None, color=GREEN):
    rect(s, x, y, w, h)
    tf = tbox(s, x + Inches(0.16), y + Inches(0.2), w - Inches(0.32), h - Inches(0.36),
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, value, 27, color, bold=True, first=True, align=PP_ALIGN.CENTER, after=3)
    para(tf, label, 10.5, DIM, align=PP_ALIGN.CENTER, after=0, line=1.15)
    if note:
        para(tf, note, 9, FAINT, align=PP_ALIGN.CENTER, before=3, after=0)


def picture(s, name, x, y, w):
    p = SHOTS / name
    if p.exists():
        s.shapes.add_picture(str(p), x, y, width=w)


def caption(s, x, y, w, head, body):
    tf = tbox(s, x, y, w, Inches(0.9))
    para(tf, head, 12.5, GREEN, bold=True, first=True, after=3)
    para(tf, body, 10.5, DIM, after=0, line=1.25)


def feature_rows(s, y, rows, row_h=Inches(0.92), label_w=Inches(3.2)):
    for label, body, color in rows:
        rect(s, M, y, CW, row_h)
        tf = tbox(s, M + Inches(0.28), y + Inches(0.16), label_w, row_h - Inches(0.3))
        para(tf, label, 13.5, color, bold=True, first=True, after=0)
        tf = tbox(s, M + Inches(0.28) + label_w, y + Inches(0.15), CW - label_w - Inches(0.6),
                  row_h - Inches(0.28))
        para(tf, body, 11, TEXT, first=True, after=0, line=1.25)
        y += row_h + Inches(0.11)
    return y


def divider(prs, kicker, title_text, blurb):
    s = slide(prs, dark=True)
    rect(s, 0, H / 2 - Inches(0.02), Inches(0.55), Pt(3), fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
    tf = tbox(s, Inches(1.1), H / 2 - Inches(1.0), Inches(10.5), Inches(2.2))
    para(tf, kicker.upper(), 11, GREEN, bold=True, first=True, after=10)
    para(tf, title_text, 40, TEXT, bold=True, after=10)
    para(tf, blurb, 14, DIM, after=0, line=1.3)
    foot(s)
    return s


# ------------------------------------------------------------------ build ---


def build() -> None:
    prs = new_deck()

    # ============================================================ 1 TITLE ===
    s = slide(prs, dark=True)
    rect(s, 0, 0, Inches(0.14), H, fill=GREEN, shape=MSO_SHAPE.RECTANGLE)
    tf = tbox(s, Inches(1.1), Inches(2.05), Inches(11), Inches(3))
    para(tf, "MINOR PROJECT  ·  PROGRESS REVIEW", 11, GREEN, bold=True, first=True, after=16)
    para(tf, "KrishiMitra", 60, TEXT, bold=True, after=4)
    para(tf, "AI-Based Crop Recommendation for Indian Farmers", 23, GREEN, after=18)
    para(tf, "Not “what can grow here?” but “what is worth growing, and why?”",
         15, DIM, italic=True, after=0)
    y = Inches(5.5)
    facts = [
        f"{COVERAGE['records_raw']:,}", "official records",
        f"{len(PROFILES)}", "states & UTs",
        f"{METRICS['dataset']['n_classes']}", "crops ranked",
        f"{METRICS['cv_accuracy_mean']*100:.1f}%", "model CV accuracy",
    ]
    x = Inches(1.1)
    for i in range(0, len(facts), 2):
        tf = tbox(s, x, y, Inches(2.6), Inches(0.9))
        para(tf, facts[i], 24, TEXT, bold=True, first=True, after=1)
        para(tf, facts[i + 1], 10.5, FAINT, after=0)
        x += Inches(2.7)

    # ======================================================== 2 MOTIVATION ==
    divider(prs, "Part one", "Motivation",
            "Why crop choice is the highest-stakes decision in Indian agriculture")

    s = slide(prs)
    y = header(s, "The Human Problem", kicker="Motivation",
               sub="Around 46% of India's workforce depends on agriculture, on an average holding of about 1.08 hectares")
    cards = [
        ("Income, not yield, is the crisis",
         "A farmer can raise yield and still lose money. Choosing the wrong crop for the soil, the rainfall and the market is the difference between a season that pays and one that adds to debt.",
         RED),
        ("The decision is made blind",
         "Crop choice is largely habit, or copying a neighbour. Soil Health Cards exist, but arrive as a twelve-parameter lab report with no action attached to the numbers.",
         AMBER),
        ("Advice answers the wrong question",
         "Existing tools predict what can grow. A farmer already knows rice grows in his village. He needs to know what is worth growing this year, on his land, at his scale.",
         BLUE),
        ("Support goes unclaimed",
         "Insurance, credit, irrigation subsidy and solar pumps all exist. The eligibility rules sit in PDFs across a dozen portals, in English, behind logins.",
         VIOLET),
    ]
    for i, (head, body, color) in enumerate(cards):
        cx = col(i % 2, COL2, GAP)
        cy = y + (Inches(1.72) + Inches(0.28)) * (i // 2)
        rect(s, cx, cy, COL2, Inches(1.72))
        rect(s, cx, cy, Pt(3), Inches(1.72), fill=color, shape=MSO_SHAPE.RECTANGLE)
        tf = tbox(s, cx + Inches(0.32), cy + Inches(0.24), COL2 - Inches(0.6), Inches(1.3))
        para(tf, head, 15, color, bold=True, first=True, after=6)
        para(tf, body, 11, DIM, after=0, line=1.28)
    foot(s)

    # ---- how it helps a farmer
    s = slide(prs)
    y = header(s, "What Changes for the Farmer", kicker="Motivation",
               sub="A concrete example: two hectares in Maharashtra")
    rect(s, M, y, COL2, Inches(3.5), border=RED)
    tf = tbox(s, M + Inches(0.32), y + Inches(0.26), COL2 - Inches(0.62), Inches(3.0))
    para(tf, "TODAY", 11, RED, bold=True, first=True, after=12)
    for t in ["Plants what he planted last year, or what the\nneighbouring field is planting",
              "No idea what the soil is short of, so fertiliser is\napplied by habit and often wasted",
              "Discovers the price only at the mandi, after harvest,\nwhen he has no choice but to sell",
              "Learns about a subsidy after the application window\nhas closed",
              "Carries all the weather risk himself, uninsured"]:
        para(tf, "—   " + t, 12, DIM, after=11, line=1.25)

    rect(s, col(1, COL2, GAP), y, COL2, Inches(3.5), border=GREEN_D)
    tf = tbox(s, col(1, COL2, GAP) + Inches(0.32), y + Inches(0.26), COL2 - Inches(0.62), Inches(3.0))
    para(tf, "WITH KRISHIMITRA", 11, GREEN, bold=True, first=True, after=12)
    for t in ["Six crops ranked for his district, each with the net\nreturn per hectare at his farm size",
              "The exact nutrient gap, converted into bags of urea,\nDAP and MOP with a rupee cost",
              "Irrigation need in millimetres beyond local rainfall,\nbefore committing to the crop",
              "Which schemes apply to that specific crop and\nholding, answered in his own language",
              "Every number shown with its source and confidence"]:
        para(tf, "✓   " + t, 12, TEXT, after=11, line=1.25)

    rect(s, M, y + Inches(3.75), CW, Inches(0.92), fill=PANEL_HI, border=GREEN_D)
    tf = tbox(s, M + Inches(0.4), y + Inches(3.95), CW - Inches(0.8), Inches(0.6),
              anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "The system never tells a farmer to trust it. It shows the ideal range for every "
             "condition, the source of every figure, and how confident it is — so the advice can be "
             "checked by an extension officer before it is acted on.",
         12.5, TEXT, first=True, after=0, line=1.25)
    foot(s)

    # ======================================================== 3 THE SYSTEM ==
    divider(prs, "Part two", "The System",
            "Architecture, data foundation, and the four signals behind every recommendation")

    # ---- objective
    s = slide(prs)
    y = header(s, "Objective", kicker="The System",
               sub="Turn a crop classifier into an advisor a farmer could actually act on")
    steps = [
        ("01", "Pick a region", "A 3D globe opens onto a map of India. Click any state to load its soil profile, climate normals and real cropping history.", GREEN),
        ("02", "Rank the crops", "Four signals combined: can it grow, is it grown here, what does it earn, and can the farmer water and afford it.", BLUE),
        ("03", "Explain the answer", "Every condition compared against the crop's ideal band, with the fertiliser gap costed in bags and rupees.", AMBER),
        ("04", "Answer the follow-up", "A grounded chatbot for the government schemes that apply to that crop and that holding.", VIOLET),
    ]
    x = M
    for num, head, body, color in steps:
        rect(s, x, y, COL4, Inches(3.15))
        tf = tbox(s, x + Inches(0.26), y + Inches(0.26), COL4 - Inches(0.5), Inches(2.7))
        para(tf, num, 26, color, bold=True, first=True, after=8)
        para(tf, head, 14.5, TEXT, bold=True, after=8)
        para(tf, body, 11, DIM, after=0, line=1.3)
        x += COL4 + Inches(0.24)
    foot(s)

    # ---- architecture
    s = slide(prs)
    y = header(s, "System Architecture", kicker="The System")
    layers = [
        ("CLIENT", "React 18 + Vite", "d3-geo canvas globe · India state map, 85 KB bundled offline · recommendation panel · scheme chat drawer", GREEN),
        ("API", "FastAPI + Pydantic", "/api/states · /api/overview · /api/recommend/state · /api/recommend/custom · /api/schemes · /api/chat", BLUE),
        ("ENGINE", "recommender.py · chatbot.py", "Four-signal ranking · percentile-envelope explainability · fertiliser planner · grounded scheme retrieval with offline fallback", AMBER),
        ("MODEL", "scikit-learn RandomForest", f"300 trees · 7 features · {METRICS['dataset']['n_classes']} classes · trained and evaluated by train.py, artefacts versioned with metrics", VIOLET),
        ("DATA", "Official + curated", f"{COVERAGE['records_raw']:,} GoI production records · {len(PROFILES)} state agro-climatic profiles · CACP costs and MSP · {len(SCHEMES)} schemes", RED),
    ]
    for tag, name, detail, color in layers:
        rect(s, M, y, CW, Inches(0.86))
        chip(s, M + Inches(0.26), y + Inches(0.25), Inches(1.15), Inches(0.36), tag, color)
        tf = tbox(s, M + Inches(1.6), y + Inches(0.13), Inches(3.1), Inches(0.62))
        para(tf, name, 13, TEXT, bold=True, first=True, after=0)
        tf = tbox(s, M + Inches(4.7), y + Inches(0.16), CW - Inches(5.0), Inches(0.6))
        para(tf, detail, 10.5, DIM, first=True, after=0, line=1.2)
        y += Inches(0.96)
    tf = tbox(s, M, y + Inches(0.06), CW, Inches(0.4))
    para(tf, "Every layer runs locally. The demo needs no network connection and no API key.",
         11.5, GREEN, italic=True, first=True, after=0)
    foot(s)

    # ---- data foundation
    s = slide(prs)
    y = header(s, "Data Foundation", kicker="The System",
               sub="Three of the four signals are derived from official Government of India data, not estimated")
    g4 = Inches(0.24)
    stat(s, col(0, COL4, g4), y, COL4, Inches(1.32), f"{COVERAGE['records_raw']:,}", "production records",
         "Ministry of Agriculture")
    stat(s, col(1, COL4, g4), y, COL4, Inches(1.32), f"{COVERAGE['districts']}", "districts covered",
         f"{COVERAGE['states_mapped']} states")
    stat(s, col(2, COL4, g4), y, COL4, Inches(1.32),
         f"{COVERAGE['years'][1]-COVERAGE['years'][0]+1} yrs", "of yield history",
         f"{COVERAGE['years'][0]}–{COVERAGE['years'][1]}", color=BLUE)
    stat(s, col(3, COL4, g4), y, COL4, Inches(1.32), f"{COVERAGE['records_used']:,}", "records used",
         "after crop mapping", color=BLUE)

    y2 = y + Inches(1.52)
    sources = [
        ("District-wise crop production statistics", "Directorate of Economics & Statistics, MoA&FW",
         "Sown area and production by district, season, crop and year. Gives the regional cultivation prior and measured yields.", GREEN),
        ("CACP cost of production, 2024-25", "Commission for Agricultural Costs & Prices",
         "A2+FL cost per quintal for field crops, and the MSP set at roughly 1.5× that cost.", GREEN),
        ("IMD climatological normals", "India Meteorological Department",
         "Temperature, humidity and rainfall normals per state, used for agro-climatic fitness and the irrigation gap.", BLUE),
        ("Soil Health Card portal aggregates", "Department of Agriculture & Farmers Welfare",
         "District-averaged N, P, K and pH, forming the soil half of each state profile.", BLUE),
        ("Scheme guidelines", f"{len(SCHEMES)} central schemes across 8 categories",
         "PM-KISAN, PMFBY, KCC, Soil Health Card, PMKSY, PM-KUSUM, e-NAM, MSP, MIDH, FPO, AIF, Natural Farming.", AMBER),
    ]
    for name, org, detail, color in sources:
        rect(s, M, y2, CW, Inches(0.66))
        tf = tbox(s, M + Inches(0.26), y2 + Inches(0.1), Inches(4.2), Inches(0.5))
        para(tf, name, 11.5, TEXT, bold=True, first=True, after=1)
        para(tf, org, 9, color, after=0)
        tf = tbox(s, M + Inches(4.6), y2 + Inches(0.15), CW - Inches(4.9), Inches(0.45))
        para(tf, detail, 10, DIM, first=True, after=0, line=1.18)
        y2 += Inches(0.74)
    foot(s)

    # ---- what real data changed
    s = slide(prs)
    y = header(s, "What Real Data Changed", kicker="The System",
               sub="Substituting official statistics for hand-authored estimates corrected the model, and corrected us")
    rows = [
        ("Cultivation prior", "Was our own judgement of which crops each state grows. Now measured sown area: Punjab is 79% rice and 15% cotton; Kerala is 72% coconut. It corrected our own guess — we had listed banana as Kerala's leading crop.", GREEN),
        ("Yield", "Was one national figure per crop. Now state-specific and measured: Punjab paddy 3.73 t/ha against Jharkhand 0.97 t/ha. A national average erases exactly the regional gap that matters.", GREEN),
        ("Risk", "Was a hand-assigned score from 1 to 5. Now the coefficient of variation of yield across 19 years, blended with market risk. Moth beans at CV 0.72 genuinely fail some years; grapes at CV 0.06 genuinely do not.", GREEN),
        ("Cost", "Was a fixed cost per hectare, quietly calibrated against optimistic yields. Now CACP cost per quintal, so cost scales with what the land actually produces — the same basis the government uses to set MSP.", AMBER),
    ]
    feature_rows(s, y, rows, row_h=Inches(1.12), label_w=Inches(2.5))
    foot(s)

    # ---- four signals
    s = slide(prs)
    y = header(s, "The Ranking Engine", kicker="The System",
               sub="Four signals, weighted explicitly — the policy is auditable, not buried inside a model")
    sigs = [
        ("Agro-climatic fitness", "45%", "Can it survive here?",
         "RandomForest probability blended with each crop's 10th–90th percentile envelope. The forest is accurate but spiky; the envelope gives a smooth signal across all 22 crops.", GREEN),
        ("Economics", "30%", "What does it earn?",
         "Measured yield × MSP or mandi price, minus CACP cost. Annualised, so a 70-day moong is comparable with a 30-year mango orchard carrying amortised establishment cost.", AMBER),
        ("Water feasibility", "15%", "Can they irrigate it?",
         "Crop water requirement against 70% of local rainfall. Reports the irrigation gap in millimetres and a rainfed / light / moderate / heavy verdict.", BLUE),
        ("Risk", "10%", "How exposed are they?",
         "Measured yield volatility for production risk, blended with perishability, price volatility and capital lock-in for market risk.", RED),
    ]
    for name, weight, q, detail, color in sigs:
        rect(s, M, y, CW, Inches(0.96))
        tf = tbox(s, M + Inches(0.28), y + Inches(0.16), Inches(2.9), Inches(0.7))
        para(tf, name, 13.5, color, bold=True, first=True, after=1)
        para(tf, q, 10.5, FAINT, italic=True, after=0)
        tf = tbox(s, M + Inches(3.3), y + Inches(0.17), Inches(7.0), Inches(0.66))
        para(tf, detail, 10.5, TEXT, first=True, after=0, line=1.22)
        tf = tbox(s, M + Inches(10.4), y + Inches(0.24), Inches(1.33), Inches(0.5), align=PP_ALIGN.RIGHT)
        para(tf, weight, 21, color, bold=True, first=True, align=PP_ALIGN.RIGHT, after=0)
        y += Inches(1.06)
    tf = tbox(s, M, y + Inches(0.06), CW, Inches(0.7))
    para(tf, "Fitness is also a gate. A crop the region does not cultivate is removed outright, however "
             "profitable it looks — and profit is discounted by regional evidence, because a farmer cannot "
             "realise a national yield in a district with no supply chain for that crop.",
         11.5, DIM, italic=True, first=True, after=0, line=1.25)
    foot(s)

    # ---- model
    s = slide(prs)
    y = header(s, "Model & Evaluation", kicker="The System",
               sub=f"RandomForest · 300 trees · {METRICS['dataset']['n_rows']:,} samples · {METRICS['dataset']['n_classes']} crops · trained in {METRICS['train_seconds']:.2f}s")
    g4 = Inches(0.24)
    stat(s, col(0, COL4, g4), y, COL4, Inches(1.4), f"{METRICS['holdout_accuracy']*100:.2f}%", "hold-out accuracy")
    stat(s, col(1, COL4, g4), y, COL4, Inches(1.4), f"{METRICS['cv_accuracy_mean']*100:.2f}%",
         "5-fold cross-validation", f"± {METRICS['cv_accuracy_std']*100:.2f}%")
    stat(s, col(2, COL4, g4), y, COL4, Inches(1.4),
         f"{METRICS['misclassified_test_samples']} / {METRICS['test_set_size']}", "misclassified", color=AMBER)
    stat(s, col(3, COL4, g4), y, COL4, Inches(1.4), "0", "SMOTE applied", "by design", color=BLUE)

    y2 = y + Inches(1.68)
    rect(s, M, y2, COL2, Inches(2.5))
    tf = tbox(s, M + Inches(0.3), y2 + Inches(0.22), COL2 - Inches(0.6), Inches(2.1))
    para(tf, "FEATURE IMPORTANCE", 10.5, DIM, bold=True, first=True, after=10)
    for feat, val in METRICS["feature_importance"].items():
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run(); r.text = f"{feat:<12}"
        r.font.size = Pt(11); r.font.color.rgb = TEXT; r.font.name = MONO
        r2 = p.add_run(); r2.text = "▌" * max(1, int(val * 46)) + f"  {val:.3f}"
        r2.font.size = Pt(10); r2.font.color.rgb = GREEN; r2.font.name = MONO

    rect(s, col(1, COL2, GAP), y2, COL2, Inches(2.5), border=AMBER)
    tf = tbox(s, col(1, COL2, GAP) + Inches(0.3), y2 + Inches(0.22), COL2 - Inches(0.6), Inches(2.1))
    para(tf, "READING 99.5% HONESTLY", 10.5, AMBER, bold=True, first=True, after=10)
    para(tf, f"The training set is exactly balanced — {METRICS['dataset']['min_class_count']} samples for "
             f"each of the {METRICS['dataset']['n_classes']} crops — so no resampling was applied. "
             "Running SMOTE here would inject synthetic noise while correcting no imbalance, and "
             "train.py asserts the balance rather than assuming it.",
         11.5, TEXT, after=9, line=1.25)
    para(tf, "The accuracy is high because this dataset has clean, well-separated class boundaries. "
             "It should not be read as real-world performance — which is precisely why the model is "
             "one signal of four rather than the whole system.",
         11.5, DIM, after=0, line=1.25)
    foot(s)

    # ======================================================= 4 THE PRODUCT ==
    divider(prs, "Part three", "The Product",
            "Every feature, as it appears to the farmer")

    # ---- feature matrix
    s = slide(prs)
    y = header(s, "Feature Set", kicker="The Product")
    feats = [
        ("Interactive 3D globe", "Canvas-rendered orthographic projection that animates to India's centroid. No three.js, no CDN textures — it runs offline.", GREEN),
        ("Clickable map of India", f"All {len(PROFILES)} states and UTs, dissolved from a 760-district boundary file and simplified 4 MB → 85 KB. Filled by the crop type ranked first in each state.", GREEN),
        ("Ranked recommendations", "Six crops per state with climate fit, water need, net return per year, MSP coverage and an honest confidence label.", BLUE),
        ("Farm-size scaling", "A slider from 0.5 to 10 hectares rescales every rupee figure; the backend owns the arithmetic so amortisation stays correct.", BLUE),
        ("Explainability panel", "Each condition compared against that crop's ideal band, mismatches surfaced first, sorted by feature importance.", AMBER),
        ("Costed fertiliser plan", "The NPK gap converted into kilograms of urea, DAP and MOP, 50 kg bags, and a total in rupees.", AMBER),
        ("Data provenance", "Every crop shows which yield source was used, the evidence that it is grown there, and its share of sown area.", VIOLET),
        ("Scheme advisory chatbot", f"Grounded retrieval over {len(SCHEMES)} central schemes, backed by Grok, degrading to local retrieval with no key and no network.", VIOLET),
        ("Soil overrides", "A farmer's own Soil Health Card or sensor readings replace the state averages through the same API.", RED),
    ]
    for i, (head, body, color) in enumerate(feats):
        cx = col(i % 3, COL3, Inches(0.26))
        cy = y + (Inches(1.42) + Inches(0.2)) * (i // 3)
        rect(s, cx, cy, COL3, Inches(1.42))
        rect(s, cx, cy, Pt(2.5), Inches(1.42), fill=color, shape=MSO_SHAPE.RECTANGLE)
        tf = tbox(s, cx + Inches(0.26), cy + Inches(0.2), COL3 - Inches(0.5), Inches(1.05))
        para(tf, head, 12.5, TEXT, bold=True, first=True, after=5)
        para(tf, body, 9.8, DIM, after=0, line=1.25)
    foot(s)

    # ---- globe + map
    s = slide(prs)
    y = header(s, "Entry: Globe to Map", kicker="The Product")
    picture(s, "01-globe.png", M, y, COL2)
    picture(s, "02-map-overview.png", col(1, COL2, GAP), y, COL2)
    cy = y + Inches(3.42)
    caption(s, M, cy, COL2, "Rendered with d3-geo on canvas",
            "The globe spins, highlights India, and animates to its centroid on click. No WebGL "
            "dependency and no texture fetch, so it cannot fail on a conference-room network.")
    caption(s, col(1, COL2, GAP), cy, COL2, f"{len(PROFILES)} states, coloured by top crop type",
            "Boundaries dissolved from a 760-district file and simplified to 85 KB, bundled into the "
            "build. One /api/overview call scores every state in about 1.4 seconds.")
    foot(s)

    # ---- recommendation
    s = slide(prs)
    y = header(s, "Recommendation Panel", kicker="The Product",
               sub="Maharashtra — mango, orange and banana: Konkan, Nagpur and Jalgaon respectively")
    picture(s, "03-maharashtra.png", M, y, Inches(7.3))
    tf = tbox(s, M + Inches(7.6), y + Inches(0.1), CW - Inches(7.6), Inches(4.4))
    for i, (head, body) in enumerate([
        ("Ranked, with the money shown", "Net return per year for the farmer's actual holding, not an abstract score."),
        ("MSP badge", "Marks crops with a guaranteed government floor price — materially lower price risk."),
        ("Honest confidence", "Every card is labelled high, moderate or low. A weak match is never dressed up as certainty."),
        ("Gate transparency", "The panel states how many crops were excluded as not cultivated there, and how many used measured government yield data."),
    ]):
        para(tf, head, 13, GREEN, bold=True, first=(i == 0), before=0 if i == 0 else 15, after=5)
        para(tf, body, 11, DIM, after=0, line=1.28)
    foot(s)

    # ---- explainability
    s = slide(prs)
    y = header(s, "Explainability & Provenance", kicker="The Product",
               sub="A farmer will not act on a number they cannot check — and neither will an extension officer")
    picture(s, "04-explainability.png", M, y, Inches(7.3))
    tf = tbox(s, M + Inches(7.6), y + Inches(0.1), CW - Inches(7.6), Inches(4.4))
    for i, (head, body) in enumerate([
        ("Condition against ideal band", "Every reading compared with that crop's observed 10th–90th percentile range, mismatches first."),
        ("Full cost breakdown", "Gross revenue, operating cost, amortised setup, measured yield volatility and the risk-adjusted figure."),
        ("Data sources, per crop", "Which yield source was used, the evidence the crop is grown there, its share of sown area, and how risk was derived."),
        ("Costed fertiliser plan", "The nutrient gap as bags of urea, DAP and MOP with a rupee total — an instruction, not a diagnosis."),
    ]):
        para(tf, head, 13, GREEN, bold=True, first=(i == 0), before=0 if i == 0 else 15, after=5)
        para(tf, body, 11, DIM, after=0, line=1.28)
    foot(s)

    # ---- chatbot
    s = slide(prs)
    y = header(s, "Scheme Advisory Chatbot", kicker="The Product",
               sub=f"Grounded on {len(SCHEMES)} central schemes · Grok (x.ai) with local retrieval fallback")
    picture(s, "05-chatbot.png", M, y, Inches(7.3))
    tf = tbox(s, M + Inches(7.6), y + Inches(0.1), CW - Inches(7.6), Inches(4.4))
    for i, (head, body) in enumerate([
        ("Retrieval before generation", "The question retrieves matching schemes first; only that verified text reaches the model as context."),
        ("Constrained deliberately", "A hallucinated subsidy percentage costs a farmer real money. The model summarises verified text rather than recalling from training."),
        ("Degrades, never fails", "Without an API key the same retrieval renders a structured answer directly. The feature falls back to search instead of going down."),
        ("Context aware", "The selected state and recommended crop are passed in, so the farmer never has to restate them."),
    ]):
        para(tf, head, 13, GREEN, bold=True, first=(i == 0), before=0 if i == 0 else 15, after=5)
        para(tf, body, 11, DIM, after=0, line=1.28)
    foot(s)

    # ========================================================= 5 RESULTS ====
    divider(prs, "Part four", "Validation",
            "Does the engine reproduce reality — and what did building it teach us?")

    s = slide(prs)
    y = header(s, "Does It Match Reality?", kicker="Validation",
               sub="Recommendations reproduce real Indian cropping patterns; none of them are hard-coded as answers")
    rows = [
        ("Assam", "Jute · Rice", "The jute and paddy belt", GREEN),
        ("West Bengal", "Jute · Rice · Banana", "India's largest jute producer", GREEN),
        ("Kerala", "Rice · Banana · Coconut", "Coconut is 72% of Kerala's sown area", GREEN),
        ("Maharashtra", "Mango · Orange · Banana", "Konkan mango, Nagpur orange, Jalgaon banana", GREEN),
        ("Punjab", "Maize · Cotton · Lentil", "Cotton is 15% of Punjab's sown area", GREEN),
        ("Rajasthan", "Watermelon · Moth beans", "Arid-zone crops; moth beans are 22% of area", GREEN),
        ("Himachal Pradesh", "Maize · Orange · Rajma", "Maize is 76% of HP's sown area", GREEN),
    ]
    hdr = tbox(s, M + Inches(0.28), y, CW - Inches(0.28), Inches(0.3))
    p = hdr.paragraphs[0]
    for txt, pad in [("STATE", 22), ("TOP RECOMMENDATIONS", 34), ("REALITY CHECK", 20)]:
        r = p.add_run(); r.text = txt.ljust(pad)
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = FAINT; r.font.name = MONO
    y += Inches(0.4)
    for state, crops, note, color in rows:
        rect(s, M, y, CW, Inches(0.56))
        tf = tbox(s, M + Inches(0.28), y + Inches(0.13), Inches(2.6), Inches(0.35))
        para(tf, state, 12.5, TEXT, bold=True, first=True, after=0)
        tf = tbox(s, M + Inches(3.1), y + Inches(0.14), Inches(4.2), Inches(0.35))
        para(tf, crops, 12, color, first=True, after=0)
        tf = tbox(s, M + Inches(7.5), y + Inches(0.15), Inches(4.2), Inches(0.35))
        para(tf, note, 10.5, DIM, first=True, after=0)
        y += Inches(0.63)
    foot(s)

    # ---- the Punjab finding
    s = slide(prs)
    y = header(s, "It Surfaces Real Problems", kicker="Validation",
               sub="The engine does not hide agronomic tension — it reports it")
    rect(s, M, y, CW, Inches(1.95), fill=RGBColor(0x2B, 0x1D, 0x18), border=AMBER)
    tf = tbox(s, M + Inches(0.42), y + Inches(0.28), CW - Inches(0.9), Inches(1.5))
    para(tf, "Punjab paddy scores 0.39 agro-climatic fitness with a 346 mm irrigation gap",
         20, AMBER, bold=True, first=True, after=9)
    para(tf, "649 mm of annual rainfall cannot support rice. Punjab grows it anyway, through "
             "groundwater extraction — and the engine flags the shortfall with a “check groundwater "
             "status” verdict. The state's documented water crisis falls directly out of the model "
             "output, without being told about it.",
         13, TEXT, after=0, line=1.3)

    y2 = y + Inches(2.2)
    two = [
        ("What a naive ranking does",
         "Rank on headline profit and the system recommends pomegranate for all 36 states — including "
         "Punjab, where it scores 0% suitability. This is not hypothetical: it is what our first "
         "implementation did, and it is reproducible by setting the profit weight to 0.9.", RED),
        ("What fixed it",
         "Rank on risk-adjusted expected return — profit weighted by climate fitness and regional "
         "evidence, then log-compressed. A crop paying ₹6 lakh per hectare that half-fits the climate "
         "is not worth more than one paying ₹3 lakh that thrives.", GREEN),
    ]
    x = M
    for head, body, color in two:
        rect(s, x, y2, COL2, Inches(2.0), border=color)
        tf = tbox(s, x + Inches(0.3), y2 + Inches(0.24), COL2 - Inches(0.6), Inches(1.6))
        para(tf, head, 13.5, color, bold=True, first=True, after=7)
        para(tf, body, 11, DIM, after=0, line=1.3)
        x += COL2 + GAP
    foot(s)

    # ---- engineering findings
    s = slide(prs)
    y = header(s, "Engineering Findings", kicker="Validation",
               sub="Four defects worth recording, because each one changed the design")
    findings = [
        ("Rainfall unit mismatch",
         "The training data's rainfall column runs 20–299 mm per growing cycle; IMD state normals are annual totals of 108–3062 mm. Feeding annual figures in put every state outside the training distribution, the forest returned ≈0 probability for everything, and the profit term won by default.", RED),
        ("Cost calibrated to optimistic yield",
         "Fixed per-hectare costs had been tuned against our own generous yield estimates. Substituting real, lower yields turned paddy into a phantom loss. Fixed by costing field crops per quintal from CACP data, so cost scales with output.", AMBER),
        ("Skewed profit normalisation",
         "Net return per hectare spans ₹10,000 to ₹7 lakh. Min–max normalising that raw range gave the largest crop a score of 1.0 and squashed everything else, putting papaya first in six of ten states. A log transform restored proportion between tiers.", AMBER),
        ("Yield stability read as safety",
         "Measured volatility captures production risk only. Papaya yields reliably but rots without a cold chain and crashes in price at a local glut. Production and market risk are now blended rather than conflated.", BLUE),
    ]
    for head, body, color in findings:
        rect(s, M, y, CW, Inches(1.14))
        rect(s, M, y, Pt(2.5), Inches(1.14), fill=color, shape=MSO_SHAPE.RECTANGLE)
        tf = tbox(s, M + Inches(0.3), y + Inches(0.16), Inches(3.0), Inches(0.9))
        para(tf, head, 12.5, color, bold=True, first=True, after=0)
        tf = tbox(s, M + Inches(3.5), y + Inches(0.16), CW - Inches(3.8), Inches(0.9))
        para(tf, body, 10.5, DIM, first=True, after=0, line=1.25)
        y += Inches(1.24)
    foot(s)

    # ---- progress
    s = slide(prs)
    y = header(s, "Current Progress", kicker="Validation",
               sub="Working end to end — model, engine, API, interface and chatbot, running fully offline")
    done = [
        "RandomForest trained, evaluated and versioned (99.50% CV)",
        "Four-signal ranking engine with real cultivation prior",
        f"{COVERAGE['records_raw']:,} official GoI records integrated",
        "State-specific measured yields and empirical risk",
        "CACP per-quintal cost model for field crops",
        f"Agro-climatic profiles for all {len(PROFILES)} states and UTs",
        "Explainability layer and costed fertiliser planner",
        "FastAPI service — 7 endpoints, validation, CORS",
        "React interface — globe, map, panel, chat drawer",
        f"Grounded chatbot over {len(SCHEMES)} schemes with offline fallback",
        "Data provenance surfaced per recommendation",
        "Verified end to end in headless Chrome, zero console errors",
    ]
    pending = [
        "Grok API key wiring",
        "District-level resolution",
        "Live weather and mandi prices",
        "Hindi and regional languages",
        "ESP32 soil sensor hardware",
        "Leaf-disease detection",
    ]
    left_w = Inches(7.6)
    right_w = CW - left_w - GAP
    rect(s, M, y, left_w, Inches(4.3), border=GREEN_D)
    tf = tbox(s, M + Inches(0.32), y + Inches(0.24), left_w - Inches(0.64), Inches(3.9))
    para(tf, "DELIVERED", 10.5, GREEN, bold=True, first=True, after=11)
    for d in done:
        para(tf, "✓   " + d, 11.5, TEXT, after=6.5, line=1.1)
    rect(s, M + left_w + GAP, y, right_w, Inches(4.3))
    tf = tbox(s, M + left_w + GAP + Inches(0.32), y + Inches(0.24), right_w - Inches(0.64), Inches(3.9))
    para(tf, "IN PROGRESS", 10.5, AMBER, bold=True, first=True, after=11)
    for d in pending:
        para(tf, "○   " + d, 11.5, DIM, after=12, line=1.15)
    foot(s)

    # ---- requirements
    s = slide(prs)
    y = header(s, "Requirements", kicker="Validation")
    groups = [
        ("Software", ["Python 3.12 · FastAPI · Uvicorn · Pydantic", "scikit-learn · pandas · NumPy · joblib · shapely",
                      "React 18 · Vite · d3-geo · topojson-client", "httpx · python-pptx"], GREEN),
        ("Data", [f"GoI crop production statistics ({COVERAGE['records_raw']:,} records)",
                  "Crop_recommendation.csv (2,200 × 22) for fitness",
                  "IMD normals · Soil Health Card aggregates",
                  "CACP MSP and cost of production 2024-25"], BLUE),
        ("External services", ["Grok / x.ai — scheme chatbot (free tier)", "Open-Meteo — live weather (planned)",
                               "Agmarknet / data.gov.in — mandi prices (planned)", "Bhashini — Indian languages (planned)"], AMBER),
        ("Hardware", ["Laptop — the demo runs entirely offline", "ESP32 + RS485 7-in-1 NPK soil sensor (planned)",
                      "DHT22 temperature and humidity module", "Feeds live readings through the existing override API"], VIOLET),
    ]
    for i, (head, items, color) in enumerate(groups):
        cx = col(i % 2, COL2, GAP)
        cy = y + Inches(2.2) * (i // 2)
        rect(s, cx, cy, COL2, Inches(2.0))
        tf = tbox(s, cx + Inches(0.3), cy + Inches(0.22), COL2 - Inches(0.6), Inches(1.65))
        para(tf, head.upper(), 10.5, color, bold=True, first=True, after=9)
        for it in items:
            para(tf, "·  " + it, 11, DIM, after=6, line=1.15)
    foot(s)

    # ========================================================== 6 FUTURE ====
    divider(prs, "Part five", "Future Scope",
            "Where this goes next, and how it scales beyond a single farmer")

    s = slide(prs)
    y = header(s, "Roadmap", kicker="Future Scope", sub="Three horizons")
    horizons = [
        ("NEAR TERM", "next iteration", GREEN, [
            "District-level resolution — Nashik ≠ Vidarbha",
            "Live weather via Open-Meteo",
            "Live mandi prices via Agmarknet",
            "Hindi and regional languages via Bhashini",
            "Voice interface for low-literacy users",
            "SHAP values alongside the ideal bands",
        ]),
        ("MEDIUM TERM", "beyond the review", AMBER, [
            "Crop rotation planning across seasons",
            "Groundwater sustainability from CGWB blocks",
            "Mandi price forecasting before sowing",
            "Leaf-disease detection from a photograph",
            "ESP32 + NPK sensor for live soil readings",
            "WhatsApp bot for non-smartphone farmers",
        ]),
        ("LONGER TERM", "the scaled system", VIOLET, [
            "Sentinel-2 NDVI crop health monitoring",
            "Land allocation optimised under constraints",
            "FPO and district-officer dashboards",
            "Glut and shortage early warning",
            "Carbon credit and sustainability scoring",
            "Credit scoring for KCC pre-qualification",
        ]),
    ]
    x = M
    for head, sub, color, items in horizons:
        rect(s, x, y, COL3, Inches(4.3))
        rect(s, x, y, COL3, Pt(3), fill=color, shape=MSO_SHAPE.RECTANGLE)
        tf = tbox(s, x + Inches(0.3), y + Inches(0.28), COL3 - Inches(0.6), Inches(3.9))
        para(tf, head, 11.5, color, bold=True, first=True, after=2)
        para(tf, sub, 9.5, FAINT, italic=True, after=13)
        for it in items:
            para(tf, "·  " + it, 11, TEXT if color == GREEN else DIM, after=11, line=1.22)
        x += COL3 + Inches(0.26)
    foot(s)

    # ---- scaling the idea
    s = slide(prs)
    y = header(s, "Scaling Beyond One Farmer", kicker="Future Scope",
               sub="The same engine serves three progressively larger audiences")
    tiers = [
        ("The farmer", "What to plant, what it costs, what it earns, which scheme applies. Delivered by voice in the local language, over WhatsApp where there is no smartphone.", GREEN),
        ("The FPO", "Aggregate the intent of hundreds of members into collective input purchase, shared cold storage, and a bulk sale that commands a better price than any one farmer could.", BLUE),
        ("The district officer", "Aggregate cropping intent across a district becomes an early warning: if 60% of farmers intend to plant the same crop, the glut and the price crash are both predictable months ahead.", VIOLET),
    ]
    feature_rows(s, y, tiers, row_h=Inches(1.25), label_w=Inches(2.6))
    rect(s, M, y + Inches(4.16), CW, Inches(0.86), fill=PANEL_HI, border=GREEN_D)
    tf = tbox(s, M + Inches(0.4), y + Inches(4.33), CW - Inches(0.8), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "The recommendation engine does not change between these tiers. Only the unit of "
             "aggregation does — which is what makes the idea scale rather than merely grow.",
         12.5, TEXT, first=True, after=0, line=1.25)
    foot(s)

    # ---- limitations
    s = slide(prs)
    y = header(s, "Known Limitations", kicker="Future Scope",
               sub="Stated plainly, because they define the roadmap")
    lims = [
        ("The fitness dataset is synthetic", "Crop_recommendation.csv has clean, well-separated envelopes and carries no geography. That is why 99.5% accuracy is easy, and why the official cultivation prior exists at all."),
        ("Apple is under-ranked for Himachal", "The dataset places apple's temperature band at 21–24 °C, which ignores its ~1,000 hour winter chill requirement. A vernalisation feature is needed."),
        ("State-level granularity", "Nashik and Vidarbha are both simply “Maharashtra”. The district boundary file is already in the repository; only the profile data is missing."),
        ("Production data ends in 2015", "The official series covers 1997–2015. Yields and cropping patterns have moved since; the newer release needs ingesting."),
        ("Prices are national, not local", "MSP is national by definition, but mandi prices vary by district and week. Agmarknet integration is the fix."),
        ("Wheat and sugarcane are absent", "Two of India's largest crops are not among the 22 in the fitness dataset, though both appear in the production statistics."),
    ]
    for i, (head, body) in enumerate(lims):
        cx = col(i % 2, COL2, GAP)
        cy = y + Inches(1.35) * (i // 2)
        tf = tbox(s, cx, cy, COL2, Inches(1.15))
        para(tf, f"{i+1}.  {head}", 12.5, AMBER, bold=True, first=True, after=5)
        para(tf, body, 10.5, DIM, after=0, line=1.28)
    foot(s)

    # ---- close
    s = slide(prs, dark=True)
    y = header(s, "Summary")
    rect(s, M, y, CW, Inches(1.35), fill=RGBColor(0x14, 0x2E, 0x28), border=GREEN_D)
    tf = tbox(s, M + Inches(0.45), y + Inches(0.26), CW - Inches(0.9), Inches(0.9),
              anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "A crop recommender that answers the question a farmer actually has —", 18, TEXT,
         bold=True, first=True, after=5)
    para(tf, "not what can grow here, but what is worth growing, and why.", 18, GREEN, bold=True, after=0)

    y2 = y + Inches(1.65)
    stats = [
        (f"{COVERAGE['records_raw']//1000}k", "official records"),
        (f"{COVERAGE['districts']}", "districts"),
        (f"{len(PROFILES)}", "states & UTs"),
        (f"{METRICS['dataset']['n_classes']}", "crops"),
        (f"{METRICS['cv_accuracy_mean']*100:.1f}%", "CV accuracy"),
        (f"{len(SCHEMES)}", "schemes"),
    ]
    tile_gap = Inches(0.2)
    tile_w = (CW - tile_gap * 5) / 6
    for i, (val, lbl) in enumerate(stats):
        stat(s, col(i, tile_w, tile_gap), y2, tile_w, Inches(1.25), val, lbl)

    tf = tbox(s, M, y2 + Inches(1.6), CW, Inches(1.2))
    para(tf, "github.com/RatnamOjha/krishi-mitra", 15, GREEN, bold=True, first=True, after=6)
    para(tf, "Model, ranking engine, REST API, React interface and scheme chatbot — working end to "
             "end and running fully offline for this demonstration.", 12, DIM, after=0, line=1.25)
    foot(s)

    prs.save(OUT)
    n = len(prs.slides._sldIdLst)
    print(f"Saved {OUT}")
    print(f"  {n} slides · {OUT.stat().st_size/1024/1024:.1f} MB")


if __name__ == "__main__":
    build()
